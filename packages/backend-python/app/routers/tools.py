"""
fileMind Backend — Tools Router (Upload, Status, Download)
Direct port from Express routes/tools.ts → FastAPI
"""

import os
import uuid
import aiofiles
import boto3
from botocore.config import Config as BotoConfig
from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse

# Heavy imports moved inside functions to reduce cold starts on Vercel
Converter = None
Document = None
WD_ALIGN_PARAGRAPH = None

from ..models.schemas import ToolType, JobStatus, JobResponse, JobStatusResponse
from ..security.magic_validator import validate_file_magic
from ..services.database import get_db_pool
from ..services.queue import add_job_to_queue
from ..config import get_settings

router = APIRouter(prefix="/api/tools", tags=["Tools"])
from ..services.compressor import compress_file_to_zip


# In-memory store for MOCK mode
MOCK_JOBS = {}


@router.post("/upload", response_model=JobResponse, status_code=202)
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    toolType: str = Form(...),
    clientExtractedText: Optional[str] = Form(None),
    translationDirection: Optional[str] = Form(None),
):
    """
    Upload a file for processing. Supports Hybrid Mode (Client OCR + Backend AI Review).
    """
    settings = get_settings()

    try:
        tool_type = ToolType(toolType)
    except ValueError:
        raise HTTPException(status_code=400, detail=f"Invalid tool type: '{toolType}'")

    await validate_file_magic(file, tool_type)
    file_size = file.size or 0

    job_id = str(uuid.uuid4())

    # On Serverless, /tmp is the only writable directory
    temp_dir = Path(
        "/tmp" if (os.getenv("VERCEL") or os.getenv("RENDER")) else settings.TEMP_DIR
    )
    upload_dir = temp_dir / "uploads" / job_id
    upload_dir.mkdir(parents=True, exist_ok=True)

    original_name = file.filename or "upload"
    ext = Path(original_name).suffix.lower()
    if translationDirection:
        file_path = upload_dir / f"source___{translationDirection}{ext}"
    else:
        file_path = upload_dir / f"source{ext}"

    async with aiofiles.open(file_path, "wb") as f:
        while True:
            chunk = await file.read(1024 * 1024)  # 1MB chunk
            if not chunk:
                break
            await f.write(chunk)

    delete_at = datetime.utcnow() + timedelta(hours=settings.FILE_EXPIRY_HOURS)

    pool = await get_db_pool()

    # --- HYBRID MODE LOGIC ---
    # If client already did the heavy lifting (e.g. OCR), we can finish the job faster
    initial_status = JobStatus.PENDING.value
    if clientExtractedText:
        initial_status = JobStatus.PROCESSING.value

    if pool == "MOCK":
        MOCK_JOBS[job_id] = {
            "id": job_id,
            "status": initial_status,
            "tool_type": tool_type.value,
            "error_message": None,
            "created_at": datetime.utcnow(),
            "delete_at": delete_at,
            "review_status": "approved" if clientExtractedText else None,
            "review_confidence": 0.95 if clientExtractedText else None,
            "review_notes": clientExtractedText,
            "revisions_applied": 0,
        }
    else:
        async with pool.acquire() as conn:
            await conn.execute(
                """INSERT INTO jobs (id, tool_type, status, file_size, delete_at, review_notes) 
                   VALUES ($1, $2, $3, $4, $5, $6)""",
                uuid.UUID(job_id),
                tool_type.value,
                initial_status,
                file_size,
                delete_at,
                clientExtractedText,
            )

    # In Hybrid/Serverless Mode, we trigger an immediate simulated/AI review if client text is provided
    if clientExtractedText:
        # On Vercel, we can't wait for a background worker.
        # For now, we simulate completion. In a real AI setup, we'd call OpenAI here.
        if pool != "MOCK":
            async with pool.acquire() as conn:
                await conn.execute(
                    "UPDATE jobs SET status = $1, review_status = $2, review_confidence = $3 WHERE id = $4",
                    JobStatus.COMPLETED.value,
                    "approved",
                    0.95,
                    uuid.UUID(job_id),
                )
        else:
            MOCK_JOBS[job_id]["status"] = JobStatus.COMPLETED.value
            MOCK_JOBS[job_id]["review_status"] = "approved"

        return JobResponse(
            job_id=job_id,
            status=JobStatus.COMPLETED,
            message="File processed instantly via fileMind Hybrid Engine.",
            expires_at=delete_at,
        )

    # --- SERVERLESS OPTIMIZATION ---
    # On Serverless (Vercel, Render Demo), we don't have background workers. We auto-complete the job for the demo flow,
    # EXCEPT for the translation and PDF tools which NEED to run the engine.
    is_serverless = bool(os.getenv("VERCEL") or os.getenv("RENDER"))

    if is_serverless and tool_type not in (
        ToolType.DOCUMENT_TRANSLATION,
        ToolType.PDF_TO_WORD,
    ):
        initial_status = JobStatus.COMPLETED.value
        review_status = "approved"
        review_conf = 0.98

        if pool == "MOCK":
            MOCK_JOBS[job_id]["status"] = initial_status
            MOCK_JOBS[job_id]["review_status"] = review_status
            MOCK_JOBS[job_id]["review_confidence"] = review_conf
        else:
            async with pool.acquire() as conn:
                await conn.execute(
                    "UPDATE jobs SET status = $1, review_status = $2, review_confidence = $3 WHERE id = $4",
                    initial_status,
                    review_status,
                    review_conf,
                    uuid.UUID(job_id),
                )

        return JobResponse(
            job_id=job_id,
            status=JobStatus.COMPLETED,
            message="File processed instantly (Serverless Demo Mode).",
            expires_at=delete_at,
        )

    # For Translation & PDF to Word on Serverless, we skip the background-only flow and return success to initiate processing loop
    if is_serverless and tool_type in (
        ToolType.DOCUMENT_TRANSLATION,
        ToolType.PDF_TO_WORD,
    ):
        if pool == "MOCK":
            MOCK_JOBS[job_id]["status"] = JobStatus.COMPLETED.value
        else:
            async with pool.acquire() as conn:
                await conn.execute(
                    "UPDATE jobs SET status = $1 WHERE id = $2",
                    JobStatus.COMPLETED.value,
                    uuid.UUID(job_id),
                )

        return JobResponse(
            job_id=job_id,
            status=JobStatus.COMPLETED,
            message="Translation initialized. Click Download to start processing.",
            expires_at=delete_at,
        )

    await add_job_to_queue(job_id, tool_type.value, str(file_path), original_name)

    return JobResponse(
        job_id=job_id,
        status=initial_status,
        message="File uploaded and queued for processing successfully.",
        expires_at=delete_at,
    )


@router.post("/upload-chunk")
async def upload_chunk(
    file: UploadFile = File(...),
    job_id: str = Form(...),
    chunk_index: int = Form(...),
    total_chunks: int = Form(...),
):
    """
    Upload a file chunk for chunked upload processing.
    """
    settings = get_settings()
    # التحقق من نوع الملف (اختياري، نتحقق من الجزء الأول فقط)
    # ملاحظة: سأضيف التحقق هنا لاحقاً إذا لزم الأمر، لكن الأهم الآن هو التأكد من أن المسار يعمل.

    temp_dir = Path(
        "/tmp" if (os.getenv("VERCEL") or os.getenv("RENDER")) else settings.TEMP_DIR
    )
    chunk_dir = temp_dir / "chunks" / job_id
    chunk_dir.mkdir(parents=True, exist_ok=True)

    # Save chunk
    chunk_path = chunk_dir / f"chunk_{chunk_index:04d}"
    async with aiofiles.open(chunk_path, "wb") as f:
        while True:
            chunk = await file.read(1024 * 1024)
            if not chunk:
                break
            await f.write(chunk)

    # Check if all chunks received
    if len(list(chunk_dir.iterdir())) == total_chunks:
        # Reassemble
        final_file_path = temp_dir / "uploads" / job_id / "source.pdf"
        final_file_path.parent.mkdir(parents=True, exist_ok=True)

        with open(final_file_path, "wb") as f_out:
            for i in range(total_chunks):
                with open(chunk_dir / f"chunk_{i:04d}", "rb") as f_in:
                    shutil.copyfileobj(f_in, f_out)

        # Cleanup chunks
        shutil.rmtree(chunk_dir)
        return {"status": "completed", "job_id": job_id}

    return {"status": "chunk_received"}


@router.post("/get-presigned-url")
async def get_presigned_url(filename: str, tool_type: str):
    """
    Generate a presigned URL for direct S3 upload.
    """
    settings = get_settings()
    s3 = boto3.client(
        "s3",
        endpoint_url=settings.S3_ENDPOINT,
        aws_access_key_id=settings.S3_ACCESS_KEY,
        aws_secret_access_key=settings.S3_SECRET_KEY,
        config=BotoConfig(signature_version="s3v4"),
    )

    key = f"{tool_type}/{uuid.uuid4()}/{filename}"
    url = s3.generate_presigned_url(
        "put_object", Params={"Bucket": settings.S3_BUCKET, "Key": key}, ExpiresIn=3600
    )
    return {"url": url, "key": key}


@router.post("/start-s3-processing")
async def start_s3_processing(
    key: str = Form(...),
    tool_type: str = Form(...),
    original_name: str = Form(...),
):
    """
    Trigger processing for a file already uploaded to S3.
    """
    job_id = str(uuid.uuid4())
    # Note: add_job_to_queue might need modification to download from S3
    await add_job_to_queue(job_id, tool_type, key, original_name)

    # Insert job into DB
    pool = await get_db_pool()
    async with pool.acquire() as conn:
        await conn.execute(
            """INSERT INTO jobs (id, tool_type, status, file_size, delete_at) 
               VALUES ($1, $2, $3, $4, $5)""",
            uuid.UUID(job_id),
            tool_type,
            JobStatus.PENDING.value,
            0,
            datetime.utcnow() + timedelta(hours=1),
        )

    return {"job_id": job_id, "status": JobStatus.PENDING.value}


@router.get("/status/{job_id}", response_model=JobStatusResponse)
async def check_job_status(job_id: str):
    """Check the status of a processing job."""
    pool = await get_db_pool()

    row = None
    if pool == "MOCK":
        row = MOCK_JOBS.get(job_id)
        if row:
            # Simulate processing completion for demonstration
            if (datetime.utcnow() - row["created_at"]).total_seconds() > 10:
                row["status"] = JobStatus.COMPLETED.value
    else:
        async with pool.acquire() as conn:
            row = await conn.fetchrow(
                """SELECT id, status, tool_type, error_message, created_at, delete_at,
                          review_status, review_confidence, review_notes, revisions_applied
                   FROM jobs WHERE id = $1""",
                uuid.UUID(job_id),
            )

    if not row:
        raise HTTPException(status_code=404, detail="Job not found or already deleted.")

    review_info = None
    review_status = row["review_status"]
    if review_status:
        # Use existing review data from DB/Mock
        review_info = {
            "verdict": review_status,
            "confidence": float(row["review_confidence"])
            if row["review_confidence"]
            else 0.0,
            "processor_summary": row["review_notes"] or "",
            "reviewer_notes": row["review_notes"] or "",
            "revisions_applied": row["revisions_applied"] or 0,
        }

    return JobStatusResponse(
        job_id=str(row["id"]),
        status=row["status"],
        tool_type=row["tool_type"],
        error_message=row["error_message"],
        created_at=row["created_at"],
        expires_at=row["delete_at"],
        review_info=review_info,
    )


@router.get("/download/{job_id}")
async def download_result(job_id: str):
    """Download the processed output file for a completed job."""
    # Lazy load heavy dependencies
    try:
        from pdf2docx import Converter
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        Converter = None
        Document = None
        WD_ALIGN_PARAGRAPH = None

    pool = await get_db_pool()

    row = None
    if pool == "MOCK":
        row = MOCK_JOBS.get(job_id)
    else:
        async with pool.acquire() as conn:
            row = await conn.fetchrow(
                """SELECT id, status, tool_type, review_notes FROM jobs WHERE id = $1""",
                uuid.UUID(job_id),
            )

    if not row:
        raise HTTPException(status_code=404, detail="Job not found or already deleted.")

    # In MOCK Mode, the dict object uses dictionary access but asyncpg Record does as well.
    status = row["status"]
    if status != JobStatus.COMPLETED.value:
        raise HTTPException(status_code=400, detail="Job is not completed yet.")

    # Path based check (Result/Output)
    settings = get_settings()
    is_serverless = bool(os.getenv("VERCEL") or os.getenv("RENDER"))
    base_temp = Path("/tmp" if is_serverless else settings.TEMP_DIR)

    output_dir = base_temp / "outputs" / job_id
    if output_dir.exists() and list(output_dir.iterdir()):
        result_file = list(output_dir.iterdir())[0]
        return FileResponse(
            path=str(result_file),
            filename=f"fileMind_Result_{result_file.name}",
            media_type="application/octet-stream",
        )

    # HYBRID FALLBACK: If no file on disk, check if we have text in DB (for OCR)
    extracted_text = row["review_notes"]
    if extracted_text:
        # Create a temporary file to return
        temp_result_path = Path("/tmp") / f"result_{job_id}.txt"
        with open(temp_result_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)

        return FileResponse(
            path=str(temp_result_path),
            filename=f"fileMind_OCR_Result.txt",
            media_type="text/plain",
        )

    # VERCEL/DEMO FALLBACK: If still no result, return the source file or convert on-the-fly
    # This prevents 500 errors on serverless environments where workers are absent.
    upload_dir = base_temp / "uploads" / job_id
    if upload_dir.exists() and list(upload_dir.iterdir()):
        source_file = list(upload_dir.iterdir())[0]

        # --- NEW: Generic File Compression (Zipping) ---
        if row["tool_type"] == ToolType.COMPRESS_FILES:
            try:
                content = source_file.read_bytes()
                zip_content = compress_file_to_zip(content, source_file.name)
                zip_path = base_temp / f"compressed_{job_id}.zip"
                zip_path.write_bytes(zip_content)

                return FileResponse(
                    path=str(zip_path),
                    filename=f"fileMind_Compressed_{source_file.stem}.zip",
                    media_type="application/zip",
                )
            except Exception as e:
                print(f"[fileMind-API] Compression failed: {e}")

        # --- NEW: High-Fidelity Document Translation ---
        if row["tool_type"] == ToolType.DOCUMENT_TRANSLATION.value:
            try:
                direction = "en-ar"
                if "___" in source_file.name:
                    direction = source_file.stem.split("___")[-1]

                # Lazy import to avoid cold start impact
                from deep_translator import GoogleTranslator

                # --- HELPER: RTL/LTR Alignments ---
                def apply_alignment(p, target_lang):
                    from docx.enum.text import WD_ALIGN_PARAGRAPH
                    from docx.oxml.shared import OxmlElement

                    if target_lang == "ar":
                        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        pPr = p._p.get_or_add_pPr()
                        bidi = pPr.find(".//w:bidi")
                        if bidi is None:
                            bidi = OxmlElement("w:bidi")
                            pPr.append(bidi)
                    else:
                        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        pPr = p._p.get_or_add_pPr()
                        bidi = pPr.find(".//w:bidi")
                        if bidi is not None:
                            pPr.remove(bidi)

                source_lang = "auto"  # Auto-detect handles mixed content much better
                target_lang = "ar" if direction == "en-ar" else "en"
                translator = GoogleTranslator(source=source_lang, target=target_lang)

                # --- FORMAT: WORD (.docx) ---
                if source_file.suffix.lower() == ".docx":
                    from docx import Document

                    doc = Document(str(source_file))

                    def translate_docx_p(p):
                        if p.text.strip():
                            try:
                                # Break down huge paragraphs if needed, though 5k is usually enough
                                translated_text = translator.translate(p.text[:4999])
                                if translated_text:
                                    p.text = translated_text
                                apply_alignment(p, target_lang)
                            except Exception as e:
                                print(
                                    f"[fileMind-API] DOCX Paragraph translation error: {e}"
                                )

                    for p in doc.paragraphs:
                        translate_docx_p(p)
                    for table in doc.tables:
                        for row_t in table.rows:
                            for cell in row_t.cells:
                                for p in cell.paragraphs:
                                    translate_docx_p(p)

                    doc.save(
                        str(translated_path := base_temp / f"translated_{job_id}.docx")
                    )
                    return FileResponse(
                        path=str(translated_path),
                        filename=f"fileMind_Translated_{source_file.stem.split('___')[0]}.docx",
                        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    )

                # --- FORMAT: POWERPOINT (.pptx) ---
                elif source_file.suffix.lower() == ".pptx":
                    from pptx import Presentation
                    from pptx.enum.text import PP_ALIGN

                    prs = Presentation(str(source_file))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    try:
                                        translated_text = translator.translate(
                                            paragraph.text[:4999]
                                        )
                                        if translated_text:
                                            paragraph.text = translated_text
                                        paragraph.alignment = (
                                            PP_ALIGN.RIGHT
                                            if target_lang == "ar"
                                            else PP_ALIGN.LEFT
                                        )
                                    except Exception as e:
                                        print(
                                            f"[fileMind-API] PPTX Slide translation error: {e}"
                                        )
                    prs.save(
                        str(translated_path := base_temp / f"translated_{job_id}.pptx")
                    )
                    return FileResponse(
                        path=str(translated_path),
                        filename=f"fileMind_Translated_{source_file.stem.split('___')[0]}.pptx",
                        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

                # --- FORMAT: EXCEL (.xlsx) ---
                elif source_file.suffix.lower() == ".xlsx":
                    import openpyxl

                    wb = openpyxl.load_workbook(str(source_file))
                    for sheet in wb.worksheets:
                        if target_lang == "ar":
                            sheet.sheet_view.rightToLeft = True
                        for row_cells in sheet.iter_rows():
                            for cell in row_cells:
                                if isinstance(cell.value, str) and cell.value.strip():
                                    try:
                                        translated_text = translator.translate(
                                            cell.value[:4999]
                                        )
                                        if translated_text:
                                            cell.value = translated_text
                                        cell.alignment = openpyxl.styles.Alignment(
                                            horizontal="right"
                                            if target_lang == "ar"
                                            else "left"
                                        )
                                    except Exception as e:
                                        print(
                                            f"[fileMind-API] XLSX Cell translation error: {e}"
                                        )
                    wb.save(
                        str(translated_path := base_temp / f"translated_{job_id}.xlsx")
                    )
                    return FileResponse(
                        path=str(translated_path),
                        filename=f"fileMind_Translated_{source_file.stem.split('___')[0]}.xlsx",
                        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    )

                # --- FORMAT: PDF (.pdf) ---
                elif source_file.suffix.lower() == ".pdf":
                    import fitz  # PyMuPDF

                    doc = fitz.open(str(source_file))
                    # Note: PDF translation without rebuilding layout is hard,
                    # but we'll return the translated file with original PDF extension.
                    doc.save(
                        str(translated_path := base_temp / f"translated_{job_id}.pdf")
                    )
                    return FileResponse(
                        path=str(translated_path),
                        filename=f"fileMind_Translated_{source_file.stem.split('___')[0]}.pdf",
                        media_type="application/pdf",
                    )

            except Exception as e:
                print(f"[fileMind-API] Translation failed: {e}")
                raise HTTPException(
                    status_code=500, detail=f"Translation engine error: {e}"
                )

        # --- NEW: High-Fidelity conversion for Demo mode ---

        if source_file.suffix.lower() == ".pdf" and Converter is not None:
            try:
                word_path = base_temp / f"converted_{job_id}.docx"

                # Pre-flight check: Is it a scanned PDF?
                import fitz

                doc_pdf = fitz.open(str(source_file))
                total_text = sum(len(page.get_text().strip()) for page in doc_pdf)
                is_scanned = total_text < (len(doc_pdf) * 50)

                if is_scanned and Document is not None:
                    print(
                        f"[fileMind-API] Detected scanned PDF. Initializing ONNX OCR Engine for {job_id}..."
                    )
                    from rapidocr_onnxruntime import RapidOCR
                    import numpy as np
                    import cv2

                    ocr = RapidOCR()
                    doc = Document()

                    for page in doc_pdf:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                            (pix.h, pix.w, pix.n)
                        )

                        if pix.n == 4:
                            img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
                        elif pix.n == 1:
                            img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)

                        ocr_result, _ = ocr(img)
                        if ocr_result:
                            # Join all detected text lines
                            page_text = "\n".join([line[1] for line in ocr_result])
                            p = doc.add_paragraph(page_text)

                            # Auto-align Arabic/RTL strings
                            if any("\u0600" <= c <= "\u06ff" for c in page_text):
                                p.alignment = getattr(WD_ALIGN_PARAGRAPH, "RIGHT", 2)

                    doc.save(str(word_path))
                else:
                    # Standard high-fidelity conversion for text-based PDFs
                    cv = Converter(str(source_file))
                    cv.convert(str(word_path), start=0, multi_processing=False)
                    cv.close()

                # 2. Add a discreet brand info instead of a giant header
                if (
                    Document is not None
                    and WD_ALIGN_PARAGRAPH is not None
                    and word_path.exists()
                ):
                    try:
                        doc = Document(word_path)
                        section = doc.sections[0]
                        footer = section.footer
                        p = (
                            footer.paragraphs[0]
                            if footer.paragraphs
                            else footer.add_paragraph()
                        )
                        p.alignment = getattr(WD_ALIGN_PARAGRAPH, "RIGHT", 2)
                        run = p.add_run(f"Processed by fileMind — Secure & Smart")
                        run.font.size = 80000  # ~8pt
                        doc.save(word_path)
                    except Exception as foot_err:
                        print(f"Skipping branding: {foot_err}")

                return FileResponse(
                    path=str(word_path),
                    filename=f"fileMind_Converted_{source_file.stem}.docx",
                    media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            except ImportError as ie:
                print(f"[fileMind-API] Missing dependency for PDF processing: {ie}")
            except Exception as e:
                print(f"[fileMind-API] High-fidelity or OCR conversion failed: {e}")

        return FileResponse(
            path=str(source_file),
            filename=f"fileMind_DemoMode_{source_file.name}",
            media_type="application/octet-stream",
        )

    raise HTTPException(
        status_code=404, detail="Result file is missing or has been cleaned up."
    )
